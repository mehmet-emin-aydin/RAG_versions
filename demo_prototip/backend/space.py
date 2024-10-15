import os
import time
import json
from typing import Dict, List, Optional
from retrying import retry
from langchain.docstore.document import Document
import requests
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_milvus import Milvus
import logging
from backend.page import Page
from transformers import GPT2Tokenizer
logger = logging.getLogger(__name__)


class Space:
    def __init__(self, confluence_client, space_name):
        self.base_url = "https://wiki.softtech.com.tr"
        self.confluence = confluence_client.confluence
        self.space_name = space_name
        
        self.embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
        # self.tokenizer = GPT2Tokenizer.from_pretrained("gpt2")
        self.connection_args = {"uri": "./milvus_demo_new.db"}
        space_folder = os.path.join('backend', self.space_name)
        # Space için klasör oluşturma
        if not os.path.exists(space_folder):
            os.makedirs(space_folder, exist_ok=True)  # exist_ok=True ekleyerek hatayı önlüyoruz

        self.page_ids_file = os.path.join(space_folder, 'page_ids.json')
        self.vectorstore = Milvus(
            self.embeddings,
            collection_name=self.space_name,
            connection_args=self.connection_args,
            consistency_level="Eventually",
            drop_old=False,
            auto_id=True,  # auto_id parametresini True yaparak kimlikleri otomatik oluşturmasını sağlıyoruz
            index_params={
                "metric_type": "COSINE",
                "index_type": "AUTOINDEX",
                "params": {}
            }
        )
    def update_space(self):
        """
        Bir space için değişen pagelerin içeriklerini ve eklerini vectorstoreda güncelleyen fonksiyon
        space klasörünün altındaki dosyadan page_ids leri alıyor. -->load_page_ids()
        space_key ile confluence apidan şu anki page_id leri ve bu page_id lerle when attributelarını alıyor. --> get_changed_ids()
        Yeni verilerle eski verileri kıyaslıyor.
        Güncel pageleri new_data dictinde tutuyor ve page_ids.json dosyasını güncelliyor. --> save_page_ids()
        Değişenleri changed_ids olarak tutuyor.
        Bu idlerle page lerin içeriklerini ve eklerini api ile alıyor. -->Page.update_page()
        Page in update edilmiş halini alıp vectorstore u upsert ediyor. --> upsert_collection()
        """
        space_start_time = time.time()
        space_size = 0
        old_data = self._load_page_ids()
        new_data, changed_ids = self._get_changed_ids(old_data, max_num_results=5000)
        # self._save_page_ids(new_data)
        # changed_ids listesini 50'şerlik parçalara bölerek işleme al
        batch_size = 50
        tokens = 0
        chunks = 0
        # for i in range(0, len(changed_ids), batch_size):
        #     # Her seferinde 50'lik bir parça oluştur
        #     up_side = min((i+50), len(changed_ids))
        #     batch = changed_ids[i: up_side]
            
        #     # updated ve created durumundaki id'leri al
        #     updated_ids = [item[0] for item in batch if item[1] in ['updated', 'created']]
            
        #     # updated_ids ile load_data fonksiyonunu çağır
        #     # documents = self.load_data(include_attachments=False, page_ids=updated_ids)
        #     # for document in documents:
        #     #     tokens += len(self.tokenizer(document.page_content)['input_ids'])
        #     # documents ve mevcut batch'i save_pages fonksiyonuna gönder
        #     [print(item[0], end = " ") for item in batch]
        #     # chunk = self.save_pages(batch, documents)
        #     # tokens += token
        #     chunks += chunk
        # # for page_id in changed_ids:
        # #     page_instance = Page(self.confluence, self.space_name, page_id)
        # #     new_size, space_docs= page_instance.get_space_content()
        # #     space_size += new_size
        space_duration = time.time() - space_start_time
        return {
            'space': self.space_name,
            'total_size_mb': space_size / (1024 * 1024),
            'duration_seconds': space_duration,
            'chunks' :chunks,
            'page_size': len(changed_ids)
        }
    def save_pages(self, changed_ids, documents):
        print("save_pages")
        from langchain.text_splitter import MarkdownHeaderTextSplitter, RecursiveCharacterTextSplitter
        from langchain.docstore.document import Document
        headers_to_split_on = [
            ("#", "Header_1"),
            ("##", "Header_2")
        ]
        markdown_splitter = MarkdownHeaderTextSplitter(headers_to_split_on=headers_to_split_on)
        

        combined_content = ""
        chunks = []
        abc = 0
        for document in documents:
            title = document.metadata['extra_info']['title']

            source = 'confluence'
            doc_id = document.metadata['doc_id']
            url = document.metadata['extra_info']['url']
            # [docs.extend(Document(page_content=chunks[i], metadata = {'source':file.name, 'page': i}) for i in range(len(chunks)))]
            metadata = {'title': title, 
                    'source': source, 
                    'doc_id': doc_id, 
                    'sub_id': 0,
                    'url': url,
                    'Header_1': '', 
                    'Header_2': ''
                    }
            markdown_document = document.page_content
            md_header_splits = markdown_splitter.split_text(markdown_document)
            
            # Recursive karakter splitter ayarları
            chunk_size = 1000
            chunk_overlap = 100
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=chunk_size, chunk_overlap=chunk_overlap
            )

        # Markdown splitlerinden oluşan belgeyi textsplitter ile tekrar böl ve sub_id ata

            for doc in md_header_splits:
                # Document'ı textsplitter ile böl
                sub_splits = text_splitter.split_documents([doc])
                
                # Eğer birden fazla parçaya bölünürse sub_id ata
                if len(sub_splits) > 1:
                    for idx, sub_doc in enumerate(sub_splits):
                        # Copy the original document's metadata and update with sub_id
                        sub_doc.metadata['sub_id'] = idx + 1
                for subsplit in sub_splits:
                    subsplit.metadata.update(metadata)
                    chunks.append(Document(page_content = subsplit.page_content, metadata = subsplit.metadata))
                    combined_content += subsplit.page_content
        print(f'chunk size: {len(chunks)}')

        # with open(f"{self.space_name}_combined.txt", "a+") as file:
        #     file.write(combined_content)
        # self.update_vector_store(changed_ids, chunks)
        # tokens=0
        # for chunk in chunks:
        #     tokens += len(self.tokenizer(chunk.page_content)['input_ids'])
    
        return len(chunks)

    def update_vector_store(self, changed_ids, documents):
        print("update_vectorstore")
        # Create a Milvus instance with auto_id set to True
        vectorstore = self.vectorstore
        
        # Collection exists, perform upsert and delete operations
        updated_ids = [item[0] for item in changed_ids if item[1] in ['updated', 'created']]
        if updated_ids:
            updated_pks = vectorstore.get_pks(f"doc_id in {updated_ids}")
            vectorstore.upsert(updated_pks, documents)  # Bu sefer ids parametresine gerek yok
    
        deleted_ids = [item[0] for item in changed_ids if item[1] == 'deleted']
        if deleted_ids:
            deleted_pks = vectorstore.get_pks(f"doc_id in {deleted_ids}")
            vectorstore.delete(deleted_pks)
        # else:
        #     # Collection does not exist, create it using from_documents
        #     vectorstore = Milvus.from_documents(
        #         documents=documents,
        #         embedding=self.embeddings,
        #         collection_name=self.space_name,
        #         connection_args=self.connection_args,
        #         consistency_level="Eventually",
        #         drop_old=True,
        #         index_params={
        #             "metric_type": "COSINE",
        #             "index_type": "AUTOINDEX",
        #             "params": {}
        #         }
        #     )


    def load_vector_store(self):

        vector_store= Milvus(
            self.embeddings,
            connection_args={"uri": "./milvus_demo_new.db"},
            collection_name=self.space_name,
        )
        return vector_store
    

    def _load_page_ids(self):
        # if os.path.exists(self.page_ids_file):
        #     with open(self.page_ids_file, 'r') as file:
        #         try:
        #             data = json.load(file)
        #             if isinstance(data, list):
        #                 return data
        #         except json.JSONDecodeError:
        #             pass
        return []


    def _save_page_ids(self, data):
        with open(self.page_ids_file, 'w') as file:
            json.dump(data, file, indent=4)


    def _get_changed_ids(self, old_data, page_status: Optional[str] = None,start = 0,max_num_results: Optional[int] = None):
        
        old_data_dict = {item['id']: item for item in old_data}
        changed_ids = []
        pages: List = []

        max_num_remaining = max_num_results
        ret = []
        while True:

            results = self.confluence.get_all_pages_from_space_raw(self.space_name, start = start, limit=(start + 50))['results']
            ret.extend(results)
            if (
                len(results) == 0
                or max_num_results is not None
                and len(results) >= max_num_remaining
            ):
                break

            start += len(results)
            if max_num_remaining is not None:
                max_num_remaining -= len(results)  


        space_pages = ret
        # space_pages = self.confluence.get_all_pages_from_space_raw(self.space_name)['results']

        for page in space_pages:
            page_id = page['id']
            page_instance = Page(self.confluence, self.space_name, page_id)
            page_info = page_instance.get_version_info()
            #onceden var olan sayfaların degisimine bakılır
            if page_id in old_data_dict:
                if old_data_dict[page_id]['when'] != page_info['when']:
                    old_data_dict[page_id]['when'] = page_info['when']
                    changed_ids.append((page_id,'updated'))
            else:
                old_data_dict[page_id] = page_info
                changed_ids.append((page_id,'created'))

        # dict1 ve dict2 içindeki id'leri set olarak oluşturma
        ids_dict1 = {item['id'] for item in old_data}
        ids_dict2 = {item['id'] for item in space_pages}
    
        # dict1'de olup dict2'de olmayan id'leri bulma
        deleted_ids = ids_dict1 - ids_dict2
        changed_ids.extend((page_id,'deleted') for page_id in deleted_ids)
        [old_data_dict.pop(id, None) for id in deleted_ids]
        return list(old_data_dict.values()), changed_ids
    



    def load_data(
        self,
        space_key: Optional[str] = None,
        page_ids: Optional[List[str]] = None,
        page_status: Optional[str] = None,
        label: Optional[str] = None,
        cql: Optional[str] = None,
        include_attachments=False,
        include_children=False,
        start: Optional[int] = None,
        cursor: Optional[str] = None,
        limit: Optional[int] = None,
        max_num_results: Optional[int] = None,
    ) -> List[Document]:
        
        import pytesseract
        pytesseract.pytesseract.tesseract_cmd = '/usr/bin/tesseract' 
        max_num_results = limit
        try:
            import html2text  # type: ignore
        except ImportError:
            raise ImportError(
                "`html2text` package not found, please run `pip install html2text`"
            )

        text_maker = html2text.HTML2Text()
        text_maker.ignore_links = True
        text_maker.ignore_images = True


        if not start:
            start = 0

        pages: List = []
        if space_key:
            pages.extend(
                self._get_data_with_paging(
                    self.confluence.get_all_pages_from_space,
                    start=start,
                    max_num_results=max_num_results,
                    space=space_key,
                    status=page_status,
                    expand="body.export_view.value",
                    content_type="page",
                )
            )
        elif page_ids:
            if include_children:
                dfs_page_ids = []
                max_num_remaining = max_num_results
                for page_id in page_ids:
                    current_dfs_page_ids = self._dfs_page_ids(
                        page_id, max_num_remaining
                    )
                    dfs_page_ids.extend(current_dfs_page_ids)
                    if max_num_results is not None:
                        max_num_remaining -= len(current_dfs_page_ids)
                        if max_num_remaining <= 0:
                            break
                page_ids = dfs_page_ids
            for page_id in (
                page_ids[:max_num_results] if max_num_results is not None else page_ids
            ):
                pages.append(
                    self._get_data_with_retry(
                        self.confluence.get_page_by_id,
                        page_id=page_id,
                        expand="body.export_view.value",
                    )
                )

        docs = []
        for page in pages:
            doc = self.process_page(page, include_attachments, text_maker)
            docs.append(doc)

        return docs
    
    def _dfs_page_ids(self, page_id, max_num_results):
        ret = [page_id]
        max_num_remaining = (
            (max_num_results - 1) if max_num_results is not None else None
        )
        if max_num_results is not None and max_num_remaining <= 0:
            return ret

        child_page_ids = self._get_data_with_paging(
            self.confluence.get_child_id_list,
            page_id=page_id,
            type="page",
            max_num_results=max_num_remaining,
        )
        for child_page_id in child_page_ids:
            dfs_ids = self._dfs_page_ids(child_page_id, max_num_remaining)
            ret.extend(dfs_ids)
            if max_num_results is not None:
                max_num_remaining -= len(dfs_ids)
                if max_num_remaining <= 0:
                    break
        return ret
    
    def _get_data_with_paging(
        self, paged_function, start=0, max_num_results=50, **kwargs
    ):
        max_num_remaining = max_num_results
        ret = []
        while True:
            results = self._get_data_with_retry(
                paged_function, start=start, limit=max_num_remaining, **kwargs
            )
            ret.extend(results)
            if (
                len(results) == 0
                or max_num_results is not None
                and len(results) >= max_num_remaining
            ):
                break

            start += len(results)
            if max_num_remaining is not None:
                max_num_remaining -= len(results)
        return ret
    
    @retry(stop_max_attempt_number=1, wait_fixed=4)
    def _get_data_with_retry(self, function, **kwargs):
        return function(**kwargs)
    



    def process_page(self, page, include_attachments, text_maker):
        logger.info("Processing " + self.base_url + page["_links"]["webui"])

        if include_attachments:
            attachment_texts = self.process_attachment(page["id"])
        else:
            attachment_texts = []
        text = text_maker.handle(page["body"]["export_view"]["value"]) + "".join(
            attachment_texts
        )
        return Document(
            page_content=text,
            metadata = {
                "doc_id": page["id"],
                "extra_info" :{
                    "title": page["title"],
                    "page_id": page["id"],
                    "status": page["status"],
                    "url": self.base_url + page["_links"]["webui"],
                }
            },
        )
    

    def process_attachment(self, page_id):
        try:
            pass
        except ImportError:
            raise ImportError(
                "`pytesseract` or `pdf2image` or `Pillow` package not found, please run"
                " `pip install pytesseract pdf2image Pillow`"
            )

        # depending on setup you may also need to set the correct path for poppler and tesseract
        attachments = self.confluence.get_attachments_from_content(page_id)["results"]
        texts = []
        for attachment in attachments:
            media_type = attachment["metadata"]["mediaType"]
            absolute_url = self.base_url + attachment["_links"]["download"]
            title = attachment["title"]
            if media_type == "application/pdf":
                logger.info("Processing PDF attachment " + absolute_url)
                text = title + self.process_pdf(absolute_url)
            elif (
                media_type == "image/png"
                or media_type == "image/jpg"
                or media_type == "image/jpeg"
                or media_type == "image/webp"
            ):
                logger.info("Processing image attachment " + absolute_url)
                text = title + self.process_image(absolute_url)
            elif (
                media_type
                == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            ):
                logger.info("Processing Word document attachment " + absolute_url)
                text = title + self.process_doc(absolute_url)
            elif (
                media_type == "application/vnd.ms-excel"
                or media_type
                == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                or media_type == "application/vnd.ms-excel.sheet.macroenabled.12"
            ):
                if title.endswith(".csv") or absolute_url.endswith(".csv"):
                    logger.info("Processing CSV attachment " + absolute_url)
                    text = title + self.process_csv(absolute_url)
                else:
                    logger.info("Processing XLS attachment " + absolute_url)
                    text = title + self.process_xls(absolute_url)
            elif media_type == "application/vnd.ms-excel.sheet.binary.macroenabled.12":
                logger.info("Processing XLSB attachment " + absolute_url)
                text = title + self.process_xlsb(absolute_url)
            elif media_type == "text/csv":
                logger.info("Processing CSV attachment " + absolute_url)
                text = title + self.process_csv(absolute_url)
            elif media_type == "application/vnd.ms-outlook":
                logger.info("Processing Outlook message attachment " + absolute_url)
                text = title + self.process_msg(absolute_url)
            elif media_type == "text/html":
                logger.info("  Processing HTML attachment " + absolute_url)
                text = title + self.process_html(absolute_url)
            elif media_type == "text/plain":
                if title.endswith(".csv") or absolute_url.endswith(".csv"):
                    logger.info("Processing CSV attachment " + absolute_url)
                    text = title + self.process_csv(absolute_url)
                else:
                    logger.info("Processing Text attachment " + absolute_url)
                    text = title + self.process_txt(absolute_url)
            elif media_type == "image/svg+xml":
                logger.info("Processing SVG attachment " + absolute_url)
                text = title + self.process_svg(absolute_url)
            elif (
                media_type
                == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                or media_type
                == "application/vnd.ms-powerpoint.presentation.macroenabled.12"
            ):
                logger.info(
                    "Processing PowerPoint attachment "
                    + absolute_url
                    + " ("
                    + media_type
                    + ")"
                )
                text = title + self.process_ppt(absolute_url)
            else:
                logger.info(
                    f"Skipping unsupported attachment {absolute_url} of media_type {media_type}"
                )
                continue
            texts.append(text)

        return texts




    def process_pdf(self, link):
        try:
            import pytesseract  # type: ignore
            from pdf2image import convert_from_bytes  # type: ignore
        except ImportError:
            raise ImportError(
                "`pytesseract` or `pdf2image` package not found, please run `pip"
                " install pytesseract pdf2image`"
            )

        response = self.confluence.request(path=link, absolute=True)
        text = ""

        if (
            response.status_code != 200
            or response.content == b""
            or response.content is None
        ):
            return text
        try:
            images = convert_from_bytes(response.content)
        except ValueError:
            return text

        for i, image in enumerate(images):
            image_text = pytesseract.image_to_string(image)
            text += f"Page {i + 1}:\n{image_text}\n\n"

        return text
    




    def process_image(self, link):
        try:
            from io import BytesIO  # type: ignore

            import pytesseract  # type: ignore
            from PIL import Image  # type: ignore
        except ImportError:
            raise ImportError(
                "`pytesseract` or `Pillow` package not found, please run `pip install"
                " pytesseract Pillow`"
            )

        text = ""

        try:
            response = self.confluence.request(path=link, absolute=True)
            # Check if the response status code indicates success (200 OK)
            if response.status_code == 200 and response.content:
                try:
                    image = Image.open(BytesIO(response.content))
                    text = pytesseract.image_to_string(image)
                except OSError:
                    # Handle errors that occur while opening or processing the image
                    logger.error(
                        f"Error processing image at {link}: Unable to open or read the image content."
                    )
                    return text
            else:
                # Log non-200 responses here if needed
                logger.error(
                    f"Error fetching image at {link}: HTTP status code {response.status_code}."
                )
                return text
        except requests.exceptions.RequestException as e:
            # This catches any Requests-related exceptions, including HTTPError, ConnectionError, etc.
            logger.error(f"Request error while fetching image at {link}: {e}")
            return text

        return text
    



    def process_doc(self, link):
        try:
            from io import BytesIO
            import docx2txt
            import zipfile  # Import zipfile to catch BadZipFile exceptions
        except ImportError:
            raise ImportError(
                "`docx2txt` package not found, please run `pip install docx2txt`"
            )

        text = ""

        try:
            response = self.confluence.request(path=link, absolute=True)
            if response.status_code != 200 or response.content in [b"", None]:
                logger.error(
                    f"Error fetching document at {link}: HTTP status code {response.status_code}."
                )
                return text

            file_data = BytesIO(response.content)
            try:
                text = docx2txt.process(file_data)
            except zipfile.BadZipFile:
                logger.error(
                    f"Error processing Word document at {link}: File is not a zip file."
                )
                return text
        except Exception as e:
            logger.error(f"Unexpected error processing document at {link}: {e}")
            return text

        return text

    def process_ppt(self, link):
        try:
            from io import BytesIO
            from pptx import Presentation  # type: ignore
        except ImportError:
            raise ImportError(
                "`python-pptx` package not found, please run `pip install python-pptx`"
            )

        response = self.confluence.request(path=link, absolute=True)
        text = ""

        if (
            response.status_code != 200
            or response.content == b""
            or response.content is None
        ):
            return text

        file_data = BytesIO(response.content)

        try:
            presentation = Presentation(file_data)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + " "
        except (
            Exception
        ) as e:  # Catching a general exception to handle any unexpected errors
            logger.error(f"Error processing PowerPoint file at {link}: {e}")
            text = f"Error processing PowerPoint file: {link}. The file might be corrupt or not a valid PowerPoint file."

        return text.strip()  # Remove any leading/trailing whitespace

    def process_xls(self, link):
        try:
            import pandas as pd  # type: ignore
        except ImportError:
            raise ImportError(
                "`pandas` package not found, please run `pip install pandas`"
            )
        try:
            from io import BytesIO
        except ImportError:
            raise ImportError("Failed to import BytesIO from io")

        response = self.confluence.request(path=link, absolute=True)
        text = ""

        if (
            response.status_code != 200
            or response.content == b""
            or response.content is None
        ):
            return text

        file_data = BytesIO(response.content)

        # Try to read the Excel file
        try:
            # Use pandas to read all sheets; returns a dict of DataFrame
            sheets = pd.read_excel(file_data, sheet_name=None, engine="openpyxl")
        except Exception as e:
            return f"Failed to read Excel file: {e!s}"

        for sheet_name, sheet_data in sheets.items():
            text += f"{sheet_name}:\n"
            for row_index, row in sheet_data.iterrows():
                text += "\t".join(str(value) for value in row) + "\n"
            text += "\n"

        return text.strip()
    




    def process_csv(self, link):
        try:
            import pandas as pd
            from io import BytesIO
        except ImportError:
            raise ImportError(
                "`pandas` package not found, please run `pip install pandas`"
            )

        response = self.confluence.request(path=link, absolute=True)
        text = ""

        if (
            response.status_code != 200
            or response.content == b""
            or response.content is None
        ):
            return text

        file_data = BytesIO(response.content)

        try:
            # Assuming CSV uses default comma delimiter. If delimiter varies, consider detecting it.
            df = pd.read_csv(file_data, low_memory=False)
            # Convert the DataFrame to a text string, including headers
            text_rows = []
            for index, row in df.iterrows():
                text_rows.append(", ".join(row.astype(str)))
            text = "\n".join(text_rows)
        except Exception as e:
            logger.error(f"Error processing CSV file: {e}")
            text = "Error processing CSV file."

        return text

    def process_svg(self, link):
        try:
            from io import BytesIO  # type: ignore

            import pytesseract  # type: ignore
            from PIL import Image  # type: ignore
            from reportlab.graphics import renderPM  # type: ignore
            from svglib.svglib import svg2rlg  # type: ignore
        except ImportError:
            raise ImportError(
                "`pytesseract`, `Pillow`, or `svglib` package not found, please run"
                " `pip install pytesseract Pillow svglib`"
            )

        response = self.confluence.request(path=link, absolute=True)
        text = ""

        if (
            response.status_code != 200
            or response.content == b""
            or response.content is None
        ):
            return text

        drawing = svg2rlg(BytesIO(response.content))

        img_data = BytesIO()
        renderPM.drawToFile(drawing, img_data, fmt="PNG")
        img_data.seek(0)
        image = Image.open(img_data)

        return pytesseract.image_to_string(image)
    


    def process_msg(self, link):
        try:
            import extract_msg  # type: ignore
            from io import BytesIO
        except ImportError:
            raise ImportError(
                "`extract-msg` package not found, please run `pip install extract-msg`"
            )

        response = self.confluence.request(path=link, absolute=True)
        text = ""

        if response.status_code != 200 or response.content in [b"", None]:
            logger.error(f"Failed to download .msg file from {link}")
            return text

        file_data = BytesIO(response.content)

        try:
            # Load the .msg file content
            with extract_msg.Message(file_data) as msg:
                subject = msg.subject
                sender = msg.sender
                to = msg.to
                cc = msg.cc
                body = msg.body

                # Compile the extracted information into a text string
                text = (
                    f"Subject: {subject}\nFrom: {sender}\nTo: {to}\nCC: {cc}\n\n{body}"
                )
        except Exception as e:
            logger.error(f"Error processing .msg file at {link}: {e}")
            return "Error processing .msg file."

        return text
    



    def process_html(self, link):
        try:
            from bs4 import BeautifulSoup  # type: ignore
            import requests
        except ImportError:
            raise ImportError(
                "`beautifulsoup4` or `requests` package not found, please run `pip install beautifulsoup4 requests`"
            )

        try:
            response = requests.get(link)
            if response.status_code != 200:
                return "Error fetching HTML content: HTTP Status Code {}".format(
                    response.status_code
                )

            # Parse the HTML content and extract text
            soup = BeautifulSoup(response.content, "html.parser")
            return soup.get_text(separator=" ", strip=True)
        except Exception as e:
            logger.error(f"Error processing HTML file at {link}: {e}")
            return f"Error processing HTML file: {link}. An error occurred while fetching or parsing the content."

    def process_txt(self, link):
        try:
            import requests
        except ImportError:
            raise ImportError(
                "`requests` package not found, please run `pip install requests`"
            )

        try:
            response = requests.get(link)
            if response.status_code != 200:
                return "Error fetching text content: HTTP Status Code {}".format(
                    response.status_code
                )
            return response.text
        except Exception as e:
            logger.error(f"Error processing text file at {link}: {e}")
            return f"Error processing text file: {link}. An error occurred while fetching the content."